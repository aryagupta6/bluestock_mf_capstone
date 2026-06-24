import os
import sqlite3
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas

# Setup paths
OUTPUT_DIR = "data/processed"
FIGURES_DIR = "reports/figures"

# Colors matching visual rules
PRIMARY_HEX = "#0A2540" # Navy
SECONDARY_HEX = "#635BFF" # Purple/Indigo
TEXT_DARK_HEX = "#1A202C" # Charcoal
ACCENT_GREEN_HEX = "#2EC4B6" # Teal
LIGHT_BG_HEX = "#F8FAFC" # Off-white
BORDER_HEX = "#E2E8F0" # Grey border

PRIMARY_RGB = RGBColor(10, 37, 64)
SECONDARY_RGB = RGBColor(99, 91, 255)
TEXT_DARK_RGB = RGBColor(26, 32, 44)
ACCENT_GREEN_RGB = RGBColor(46, 196, 182)
LIGHT_BG_RGB = RGBColor(248, 250, 252)

class NumberedCanvas(canvas.Canvas):
    """
    Custom canvas to handle two-pass page numbering ('Page X of Y')
    and professional headers/footers.
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_page_decorations(self, page_count):
        self.saveState()
        
        # Cover page (Page 1) has its own style, skip header/footer decoration
        if self._pageNumber == 1:
            self.restoreState()
            return
            
        # Draw Running Header
        self.setFont("Helvetica-Bold", 8)
        self.setFillColor(colors.HexColor(PRIMARY_HEX))
        self.drawString(54, 755, "BLUESTOCK MUTUAL FUND CAPSTONE PROJECT")
        self.setFont("Helvetica", 8)
        self.setFillColor(colors.HexColor("#4A5568"))
        self.drawRightString(558, 755, "PORTFOLIO PERFORMANCE & RISK ANALYTICS")
        
        self.setStrokeColor(colors.HexColor("#CBD5E1"))
        self.setLineWidth(0.5)
        self.line(54, 748, 558, 748)
        
        # Draw Running Footer
        self.line(54, 52, 558, 52)
        self.drawString(54, 38, "Confidential - For Internal Review Only")
        self.drawCentredString(306, 38, "Bluestock Fintech")
        
        page_str = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 38, page_str)
        
        self.restoreState()

def add_slide_header(slide, title_text):
    """Adds a standard header to a content slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Georgia"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_RGB

def format_bullet_point(tf, text, level=0, bold_prefix=None, space_after=10):
    """Formats a paragraph as a bullet point."""
    p = tf.add_paragraph() if len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
    p.level = level
    p.space_after = Pt(space_after)
    
    if bold_prefix:
        run_bold = p.add_run()
        run_bold.text = bold_prefix
        run_bold.font.name = "Calibri"
        run_bold.font.size = Pt(14)
        run_bold.font.bold = True
        run_bold.font.color.rgb = TEXT_DARK_RGB
        
    run_text = p.add_run()
    run_text.text = text
    run_text.font.name = "Calibri"
    run_text.font.size = Pt(14)
    run_text.font.color.rgb = TEXT_DARK_RGB

def build_presentation(output_path):
    print("[*] Generating PowerPoint Presentation...")
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title Slide (Dark Theme)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape for dark theme
    bg = slide.shapes.add_shape(
        1, # Rectangle
        0, 0, Inches(10), Inches(5.625)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_RGB
    bg.line.fill.background() # No border
    
    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Mutual Fund Portfolio Analytics"
    p.font.name = "Georgia"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "End-to-End Performance Scorecard, Advanced Risk Metrics & Delivery"
    p2.font.name = "Calibri"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_GREEN_RGB
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Prepared by: Data Analyst Team | Bluestock Capstone Project"
    p3.font.name = "Calibri"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(200, 200, 200)

    # Slide 2: Objectives & Context
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "1. Project Scope & Key Objectives")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Develop a robust, relational data storage layer in SQLite using a star schema design.", bold_prefix="Relational Database Layer: ")
    format_bullet_point(tf, "Construct a programmatic ETL pipeline to clean, ingest, and integrate transaction logs, NAV history, and fund master files.", bold_prefix="Automated ETL Pipeline: ")
    format_bullet_point(tf, "Calculate industry-standard performance metrics including CAGR (1y/3y/5y), Sharpe/Sortino ratios, and OLS Alpha/Beta.", bold_prefix="Performance Engineering: ")
    format_bullet_point(tf, "Perform advanced analytics on Value at Risk (VaR), sector concentration (HHI), and investor behavioral cohorts.", bold_prefix="Advanced Risk Analytics: ")
    format_bullet_point(tf, "Deliver interactive reporting interfaces (Power BI) and automated delivery files.", bold_prefix="Production Delivery: ")

    # Slide 3: Relational DB Schema
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "2. Relational Database & Star Schema Design")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Organized around a central Dimension-Fact structure to support efficient analytical queries.", bold_prefix="Design Philosophy: ")
    format_bullet_point(tf, "dim_fund (AMFI code PK, fund house, category, fees, risk profile).", bold_prefix="Dimension Tables: ")
    format_bullet_point(tf, "fact_nav (historical NAVs, computed daily returns), fact_transactions (investor logs), and fact_performance.", bold_prefix="Fact Tables: ")
    format_bullet_point(tf, "Applied indexes on nav_date, amfi_code, and transaction_date to speed up time-series lookups.", bold_prefix="Indexing Strategy: ")
    
    # Graphic box representing tables
    shape = slide.shapes.add_shape(1, Inches(5.3), Inches(1.4), Inches(4.2), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG_RGB
    shape.line.color.rgb = SECONDARY_RGB
    stf = shape.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "SQLite Data Model Schema"
    sp.font.name = "Georgia"
    sp.font.bold = True
    sp.font.size = Pt(14)
    sp.font.color.rgb = PRIMARY_RGB
    sp.space_after = Pt(8)
    
    tables_list = [
        "- dim_fund (Fund Master)",
        "- fact_nav (NAV history & returns)",
        "- fact_transactions (Investor ledger)",
        "- fact_performance (Scorecard metrics)",
        "- fact_portfolio_holdings (Stock allocations)",
        "- fact_benchmark (Index closing history)",
        "- fact_aum & fact_sip_inflows (Industry metrics)"
    ]
    for t in tables_list:
        p = stf.add_paragraph()
        p.text = t
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK_RGB
        p.space_after = Pt(4)

    # Slide 4: Data Cleaning & Inflow Auditing
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "3. Data Cleaning, Imputation & Auditing")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Imputed missing daily NAV data using forward-fill (FFill) grouped by AMFI code to preserve market time series without look-ahead bias.", bold_prefix="Missing Value Imputation: ")
    format_bullet_point(tf, "Audited and removed invalid negative transaction amounts in the investor ledger to secure mathematical integrity.", bold_prefix="Ledger Validation: ")
    format_bullet_point(tf, "Standardized schema names, transaction types (SIP, Lumpsum, Redemption), and KYC status columns.", bold_prefix="Categorical Standardisation: ")
    format_bullet_point(tf, "Cross-referenced daily fund returns against benchmark index returns to detect anomalous movements.", bold_prefix="Quality Checks: ")

    # Slide 5: EDA Highlights
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "4. Exploratory Data Analysis (EDA) Insights")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "AUM distribution shows high concentration in top 10 AMC houses (such as HDFC, SBI, ICICI), leaving a long tail of smaller funds.", bold_prefix="AUM Concentration: ")
    format_bullet_point(tf, "Regular plans charge significantly higher expense ratios than direct plans (averaging 0.8% - 1.2% higher per year).", bold_prefix="Plan Cost Spread: ")
    format_bullet_point(tf, "Equity funds outperform debt funds on absolute returns but experience substantially higher annualized standard deviation.", bold_prefix="Asset Class Performance: ")
    format_bullet_point(tf, "Monthly SIP inflows grew steadily across the 2022-2025 period, indicating resilient retail investor participation.", bold_prefix="SIP Growth Trend: ")

    # Slide 6: Performance Scorecard Method
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "5. The Weighted Fund Scorecard Model")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Ensures a balanced, multi-dimensional view of performance rather than looking solely at raw historical returns.", bold_prefix="Objective: ")
    format_bullet_point(tf, "CAGR 3-Year (30% weight) and Sharpe Ratio (25% weight).", bold_prefix="Return & Risk-Adjusted: ")
    format_bullet_point(tf, "OLS Benchmark Alpha (20% weight) and Maximum Drawdown (10% weight, inverse).", bold_prefix="Manager Alpha & Tail Risk: ")
    format_bullet_point(tf, "Expense Ratio (15% weight, inverse) to penalize high management fees.", bold_prefix="Cost Efficiency: ")
    format_bullet_point(tf, "Scores are computed using percentile ranks (0 to 100 scale) within categories to ensure fair comparisons.", bold_prefix="Normalisation Method: ")

    # Slide 7: Top & Bottom Performers
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "6. Top & Bottom Performers by Category")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "HDFC Mid-Cap Opportunities Fund (Score: 81.1)", bold_prefix="Equity Top: ")
    format_bullet_point(tf, "ABSL Frontline Equity Fund (Score: 69.4)", bold_prefix="Equity Second: ")
    format_bullet_point(tf, "UTI Flexi Cap Fund (Score: 50.1)", bold_prefix="Equity Third: ")
    format_bullet_point(tf, "Axis Bluechip Fund - Regular (Score: 38.2)", bold_prefix="Equity Bottom: ")
    
    tb2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(3.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    format_bullet_point(tf2, "Kotak Liquid Fund - Regular (Score: 39.5)", bold_prefix="Debt Top: ")
    format_bullet_point(tf2, "SBI Magnum Gilt Fund - Regular (Score: 36.9)", bold_prefix="Debt Second: ")
    format_bullet_point(tf2, "Nippon India Gilt Securities Fund (Score: 35.9)", bold_prefix="Debt Third: ")
    format_bullet_point(tf2, "ABSL Short Term Debt Fund (Score: 28.5)", bold_prefix="Debt Bottom: ")

    # Slide 8: Advanced Risk Metrics (VaR/CVaR)
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "7. Tail Risk: Value at Risk (VaR) & CVaR")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Historical VaR (95%) represents the daily loss threshold that will only be breached 5% of the time.", bold_prefix="Value at Risk (95%): ")
    format_bullet_point(tf, "Conditional VaR (CVaR) represents the average loss expected on those worst 5% of days.", bold_prefix="Conditional VaR (95%): ")
    format_bullet_point(tf, "Mid cap and Small cap funds show a daily VaR 95% of -1.8% to -2.3%, indicating high short-term drawdown risk.", bold_prefix="Equity Risk Profile: ")
    format_bullet_point(tf, "Liquid and Short Term Debt funds present a VaR of -0.01% to -0.15%, making them ideal for capital preservation.", bold_prefix="Debt Risk Profile: ")

    # Slide 9: Rolling 90-Day Sharpe Ratio
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "8. Rolling 90-Day Sharpe Ratio Trends")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Evaluates how risk-adjusted returns evolve over market cycles, highlighting return consistency.", bold_prefix="Rolling Sharpe (90 Days): ")
    format_bullet_point(tf, "Liquid and Short Term Debt funds show very high and stable rolling Sharpe ratios due to extremely low volatility.", bold_prefix="Liquid Fund Stability: ")
    format_bullet_point(tf, "Equity fund Sharpe ratios fluctuate wildly, ranging from negative during market corrections to over 2.5 during rallies.", bold_prefix="Equity Volatility: ")
    
    # Embed the pre-computed rolling Sharpe chart
    chart_path = os.path.join(FIGURES_DIR, "rolling_sharpe_chart.png")
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(5.2), Inches(1.3), Inches(4.4), Inches(3.5))
    else:
        # Placeholder shape
        shape = slide.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.4), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG_RGB

    # Slide 10: Cohorts & Behavioral Analytics
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "9. Investor Cohort & Behavioral Analytics")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "2024 cohort (1362 investors) vs. 2025 cohort (12 investors) shows transaction volume concentration in 2024.", bold_prefix="Acquisition Cohorts: ")
    format_bullet_point(tf, "Average SIP size remains highly stable across cohorts (around Rs. 4,500 - 5,000 per month).", bold_prefix="SIPSizing: ")
    format_bullet_point(tf, "Out of 1362 active SIP clients, 1332 (97.8%) were flagged as 'at-risk' due to transaction gaps exceeding 35 days.", bold_prefix="SIP Continuity: ")
    format_bullet_point(tf, "High attrition highlight the need for automated reminders and mandate tracking to prevent SIP defaults.", bold_prefix="Retention Insight: ")

    # Slide 11: Recommender Engine & HHI
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "10. Sector HHI & Recommender Logic")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Assesses sector concentration. HHI > 2500 represents highly concentrated thematic portfolios.", bold_prefix="Sector HHI: ")
    format_bullet_point(tf, "The fund recommendation engine is built as an interactive command-line interface.", bold_prefix="CLI Recommender: ")
    format_bullet_point(tf, "Maps risk profile (Low, Moderate, High) and Asset Class (Equity, Debt) to select top 3 funds sorted by Sharpe ratio.", bold_prefix="Recommender Rules: ")
    
    chart_path_hhi = os.path.join(FIGURES_DIR, "sector_hhi_chart.png")
    if os.path.exists(chart_path_hhi):
        slide.shapes.add_picture(chart_path_hhi, Inches(5.2), Inches(1.3), Inches(4.4), Inches(3.5))
    else:
        shape = slide.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.4), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG_RGB

    # Slide 12: Summary & Recommendations
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "11. Summary & Actionable Recommendations")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    format_bullet_point(tf, "Implement automated payment mandate checks to address the 97.8% default risk flagged in the SIP continuity analysis.", bold_prefix="For Operations: ")
    format_bullet_point(tf, "Promote Direct Plans to yield-sensitive clients; regular plans impose an unnecessary drag on long-term compound growth.", bold_prefix="For Advisory: ")
    format_bullet_point(tf, "Use the weighted fund scorecard instead of raw returns to benchmark fund managers and detect performance decay early.", bold_prefix="For Performance: ")
    format_bullet_point(tf, "Deploy the recommender engine as a microservice to assist advisory teams in offering objective, risk-aligned recommendations.", bold_prefix="For Tech: ")

    prs.save(output_path)
    print(f"[+] PowerPoint presentation compiled successfully at {output_path}")

def draw_cover_page(canvas_obj, doc):
    canvas_obj.saveState()
    # Draw dark background on cover page
    canvas_obj.setFillColor(colors.HexColor(PRIMARY_HEX))
    canvas_obj.rect(0, 0, 612, 792, fill=True, stroke=False)
    
    # Title
    canvas_obj.setFillColor(colors.white)
    canvas_obj.setFont("Helvetica-Bold", 32)
    canvas_obj.drawString(54, 450, "Mutual Fund Portfolio")
    canvas_obj.drawString(54, 400, "Performance & Risk Analytics")
    
    # Subtitle
    canvas_obj.setFillColor(colors.HexColor(ACCENT_GREEN_HEX))
    canvas_obj.setFont("Helvetica-Bold", 14)
    canvas_obj.drawString(54, 350, "Capstone Project - Final Delivery Report")
    
    # Info
    canvas_obj.setFillColor(colors.HexColor("#CBD5E1"))
    canvas_obj.setFont("Helvetica", 10)
    canvas_obj.drawString(54, 150, "PREPARED FOR: Bluestock Fintech")
    canvas_obj.drawString(54, 130, "PREPARED BY: Data Analyst Team")
    canvas_obj.drawString(54, 110, "DATE: June 2026")
    
    canvas_obj.restoreState()

def build_pdf_report(output_path):
    print("[*] Generating PDF Report...")
    
    # Setup document
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=54,
        rightMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=colors.HexColor(PRIMARY_HEX),
        spaceAfter=15
    )
    
    h1_style = ParagraphStyle(
        'SectionH1',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor(PRIMARY_HEX),
        spaceBefore=15,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'SubSectionH2',
        parent=styles['Heading3'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=15,
        textColor=colors.HexColor(SECONDARY_HEX),
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'ReportBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor(TEXT_DARK_HEX),
        spaceAfter=8
    )
    
    bullet_style = ParagraphStyle(
        'ReportBullet',
        parent=body_style,
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=6
    )
    
    table_text_style = ParagraphStyle(
        'TableText',
        parent=body_style,
        fontSize=8,
        leading=10,
        spaceAfter=0
    )
    
    table_header_style = ParagraphStyle(
        'TableHeader',
        parent=table_text_style,
        fontName='Helvetica-Bold',
        textColor=colors.white
    )

    story = []
    
    # 1. Cover Page
    story.append(Spacer(1, 100)) # Filler for cover layout
    story.append(PageBreak())
    
    # 2. Executive Summary
    story.append(Paragraph("Executive Summary", title_style))
    story.append(Paragraph(
        "This capstone project presents an end-to-end investment analytics platform designed to evaluate "
        "mutual fund performance and monitor portfolio risk. By building an automated data engineering pipeline, "
        "we cleaned, integrated, and loaded disparate industry data files (NAV history, transaction logs, stock holdings, "
        "and fund attributes) into a centralized relational SQLite database.",
        body_style
    ))
    story.append(Paragraph(
        "Using this database, we implemented a composite weighted scorecard model that ranks funds based on "
        "risk-adjusted returns (Sharpe ratio), historical consistency (3y CAGR), manager skill (benchmark alpha), "
        "expense loads, and maximum drawdowns. Additionally, we executed advanced tail-risk calculations (VaR/CVaR), "
        "investor cohort behavioral analytics, and portfolio concentration assessments. Finally, an interactive recommender "
        "engine was deployed to dynamically match user risk appetites with optimal investment portfolios.",
        body_style
    ))
    story.append(Spacer(1, 15))
    
    # 3. Relational Schema Design
    story.append(Paragraph("1. Relational Database & ETL Pipeline", h1_style))
    story.append(Paragraph(
        "To support high-performance analytical queries, a relational star schema database was implemented in SQLite. "
        "The schema isolates dimension data (dim_fund) from transactional and time-series fact data (fact_nav, fact_transactions, "
        "fact_portfolio_holdings, fact_performance, and fact_benchmark). This structure minimizes data redundancy and optimizes join speeds.",
        body_style
    ))
    
    # SQLite tables summary table
    db_summary_data = [
        [Paragraph("Table Name", table_header_style), Paragraph("Type", table_header_style), Paragraph("Description", table_header_style)],
        [Paragraph("dim_fund", table_text_style), Paragraph("Dimension", table_text_style), Paragraph("Fund master sheet containing fees, plans, launch dates, and fund managers.", table_text_style)],
        [Paragraph("fact_nav", table_text_style), Paragraph("Fact", table_text_style), Paragraph("Daily Net Asset Value (NAV) time series with pre-calculated daily returns.", table_text_style)],
        [Paragraph("fact_transactions", table_text_style), Paragraph("Fact", table_text_style), Paragraph("Ledger of investor transaction details including demographic information.", table_text_style)],
        [Paragraph("fact_performance", table_text_style), Paragraph("Fact", table_text_style), Paragraph("Performance analytics including Sharpe, Sortino, CAGR, and Drawdowns.", table_text_style)],
        [Paragraph("fact_portfolio_holdings", table_text_style), Paragraph("Fact", table_text_style), Paragraph("Underlying stock allocations and weights per equity fund.", table_text_style)],
        [Paragraph("fact_benchmark", table_text_style), Paragraph("Fact", table_text_style), Paragraph("Historical closing price time series of benchmark indices.", table_text_style)]
    ]
    t_db = Table(db_summary_data, colWidths=[100, 70, 334])
    t_db.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor(BORDER_HEX)),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor(LIGHT_BG_HEX)])
    ]))
    story.append(t_db)
    story.append(Spacer(1, 15))
    
    # 4. Data Cleaning & Validation Auditing
    story.append(Paragraph("2. Data Cleaning & Quality Control", h1_style))
    story.append(Paragraph(
        "A critical phase of the data pipeline involved cleaning and sanitizing raw source files before database insertion:",
        body_style
    ))
    story.append(Paragraph("<b>- NAV Time-Series FFill:</b> Missing daily NAV records were imputed using forward-fill (FFill) grouped by AMFI code, ensuring that weekends or non-trading days did not create gaps in rolling metrics.", bullet_style))
    story.append(Paragraph("<b>- Transaction Ledger Validation:</b> Audited and removed rows containing invalid non-positive transaction amounts (amount <= 0). Adjusted text representations to standard categorical representations (SIP, Lumpsum, Redemption).", bullet_style))
    story.append(Paragraph("<b>- Plan and Option Classification:</b> Standardized plans (Regular vs. Direct) and options (Growth vs. Dividend) using regular expressions to resolve AMC-specific nomenclature naming errors.", bullet_style))
    story.append(Spacer(1, 15))
    
    # 5. Performance Engineering
    story.append(Paragraph("3. Performance & Risk-Adjusted Return Metrics", h1_style))
    story.append(Paragraph(
        "Using the clean historical NAV series, we calculated standard performance metrics for all funds:",
        body_style
    ))
    story.append(Paragraph("<b>- CAGR (1y/3y/5y):</b> Compound annual growth rates calculated peak-to-peak using actual trading days between target start and end dates.", bullet_style))
    story.append(Paragraph("<b>- Sharpe and Sortino Ratios:</b> Computed daily risk-adjusted ratios and annualized them using a 252-day multiplier. The risk-free rate was set at 6.5% annually. Sortino downside deviation was calculated using negative daily returns only.", bullet_style))
    story.append(Paragraph("<b>- Alpha & Beta (OLS):</b> Executed OLS linear regression of daily fund returns against the Nifty 100 benchmark (index_name = 'NIFTY100') to obtain beta (regression slope) and annualized alpha (intercept scaled by 252).", bullet_style))
    
    # Top 5 scorecard table
    story.append(Paragraph("Top Performing Funds (Overall Weighted Scorecard)", h2_style))
    scorecard_path = "fund_scorecard.csv"
    if os.path.exists(scorecard_path):
        df_sc = pd.read_csv(scorecard_path).sort_values('overall_rank').head(5)
        sc_data = [
            [Paragraph("Rank", table_header_style), Paragraph("Fund Name", table_header_style), Paragraph("Category", table_header_style), Paragraph("Sharpe", table_header_style), Paragraph("3y CAGR", table_header_style), Paragraph("Score", table_header_style)]
        ]
        for _, r in df_sc.iterrows():
            sc_data.append([
                Paragraph(str(int(r['overall_rank'])), table_text_style),
                Paragraph(str(r['scheme_name']).split(" - ")[0], table_text_style),
                Paragraph(str(r['category']), table_text_style),
                Paragraph(f"{r['sharpe_ratio']:.2f}", table_text_style),
                Paragraph(f"{r['cagr_3yr']*100:.1f}%", table_text_style),
                Paragraph(f"{r['weighted_score']:.1f}", table_text_style)
            ])
        t_sc = Table(sc_data, colWidths=[35, 230, 80, 50, 50, 59])
        t_sc.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
            ('TOPPADDING', (0,0), (-1,-1), 5),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor(BORDER_HEX)),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor(LIGHT_BG_HEX)])
        ]))
        story.append(t_sc)
    story.append(Spacer(1, 15))

    story.append(PageBreak())

    # 6. Advanced Analytics
    story.append(Paragraph("4. Advanced Risk & Behavioral Analytics", h1_style))
    story.append(Paragraph(
        "To deepen our risk assessment, advanced metrics were computed across the database:",
        body_style
    ))
    
    story.append(Paragraph("4.1 Value at Risk (VaR) & Conditional VaR (CVaR)", h2_style))
    story.append(Paragraph(
        "Value at Risk (95% daily VaR) identifies the threshold loss that is exceeded only 5% of the time. "
        "Conditional VaR (CVaR) represents the expected average return on those worst 5% of trading days. "
        "Equity funds exhibited substantial downside tails (VaR ranging from -1.8% to -2.3% daily), whereas debt "
        "and liquid funds maintained extremely tight profiles (VaR < -0.15% daily), validating their role as wealth-preservation vehicles.",
        body_style
    ))
    
    story.append(Paragraph("4.2 Sector Concentration (HHI) for Equity Funds", h2_style))
    story.append(Paragraph(
        "The Herfindahl-Hirschman Index (HHI) was calculated across all equity funds based on their underlying "
        "sector weights. Equity portfolios with an HHI below 1500 are highly diversified, while funds exceeding 2500 "
        "exhibit significant concentration (sector-focused portfolios). This allows advisors to flag concentration risk in thematic portfolios.",
        body_style
    ))
    
    story.append(Paragraph("4.3 Investor Behavioral Cohort Analysis", h2_style))
    cohort_path = "data/processed/cohort_analysis.csv"
    if os.path.exists(cohort_path):
        df_co = pd.read_csv(cohort_path)
        co_data = [
            [Paragraph("Cohort Year", table_header_style), Paragraph("Investors", table_header_style), Paragraph("Avg SIP (INR)", table_header_style), Paragraph("Total Invested (Cr)", table_header_style), Paragraph("Preferred Fund", table_header_style)]
        ]
        for _, r in df_co.iterrows():
            co_data.append([
                Paragraph(str(int(r['cohort_year'])), table_text_style),
                Paragraph(f"{int(r['investor_count']):,}", table_text_style),
                Paragraph(f"Rs. {r['avg_sip_size_inr']:,.2f}", table_text_style),
                Paragraph(f"Rs. {r['total_invested_inr']/10000000:.2f} Cr", table_text_style),
                Paragraph(str(r['top_preferred_fund']).split(" - ")[0], table_text_style)
            ])
        t_co = Table(co_data, colWidths=[70, 70, 90, 100, 174])
        t_co.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
            ('TOPPADDING', (0,0), (-1,-1), 5),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor(BORDER_HEX)),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor(LIGHT_BG_HEX)])
        ]))
        story.append(t_co)
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("4.4 SIP Continuity & At-Risk Accounts", h2_style))
    story.append(Paragraph(
        "For clients with 6 or more monthly SIP transactions, we tracked the average interval (gap) between consecutive "
        "payment dates. If an investor's average gap exceeds 35 days, they are flagged as an 'at-risk' account showing "
        "signs of payment default. Out of 1,362 eligible SIP clients, 1,332 (97.8%) were flagged as at-risk. This high rate "
        "underlines a critical operational vulnerability where payment mandates are either failing or not being consistently renewed.",
        body_style
    ))
    
    # 7. Strategic Recommendations
    story.append(Spacer(1, 15))
    story.append(Paragraph("5. Strategic Advisory & Operational Recommendations", h1_style))
    story.append(Paragraph("<b>- Address SIP Mandate Defaults:</b> Operational teams should immediately establish automated email and SMS alerts 5 days prior to SIP dates to reduce the 97.8% default risk.", bullet_style))
    story.append(Paragraph("<b>- Push Direct Plans:</b> Financial advisors must highlight the compounding benefit of Direct Plans. Direct plans skip broker commissions, generating significantly higher Sharpe ratios and long-term AUM.", bullet_style))
    story.append(Paragraph("<b>- Use Scorecard Benchmarking:</b> Fund research teams should replace simple return-ranking systems with the composite scorecard, capturing tail risk (Drawdown) and cost (Expense) alongside absolute performance.", bullet_style))
    
    # Build document
    doc.build(story, canvasmaker=NumberedCanvas, onFirstPage=draw_cover_page)
    print(f"[+] PDF report generated successfully at {output_path}")

def main():
    print("="*60)
    print("FINAL DELIVERY GENERATOR (DAY 7)")
    print("="*60)
    
    os.makedirs("reports", exist_ok=True)
    
    # Compile PPTX
    pptx_path = "reports/Presentation.pptx"
    build_presentation(pptx_path)
    
    # Compile PDF
    pdf_path = "reports/Final_Report.pdf"
    build_pdf_report(pdf_path)
    
    print("="*60)
    print("[+] All final delivery files compiled successfully!")
    print("="*60)

if __name__ == "__main__":
    main()
